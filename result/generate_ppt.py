#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""铝电解槽工作电压参数优化项目 — PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Arc
import numpy as np
import os
import io

# ── 全局配色 ──────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
STEEL     = RGBColor(0x2E, 0x50, 0x90)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY  = RGBColor(0x66, 0x66, 0x66)
ACCENT_R  = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_G  = RGBColor(0x27, 0xAE, 0x60)
ACCENT_O  = RGBColor(0xE6, 0x7E, 0x22)
LIGHT_LINE= RGBColor(0xDC, 0xE3, 0xEE)
BG_GRAD   = RGBColor(0xEE, 0xF2, 0xF8)

# Matplotlib 配色
MPL_NAVY   = '#1B2A4A'
MPL_STEEL  = '#2E5090'
MPL_ACCENT = '#3A6EA5'
MPL_RED    = '#C0392B'
MPL_GREEN  = '#27AE60'
MPL_ORANGE = '#E67E22'
MPL_LIGHT  = '#E8F0FE'
MPL_GRAY     = '#95A5A6'
MPL_DARK_GRAY = '#2D2D2D'
MPL_MID_GRAY  = '#666666'

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
CHARTS_DIR = os.path.join(OUTPUT_DIR, 'ppt_charts')
os.makedirs(CHARTS_DIR, exist_ok=True)

# Matplotlib 中文字体
plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DengXian', 'KaiTi']
plt.rcParams['axes.unicode_minus'] = False


# ── 图表生成函数 ────────────────────────────────────────

def chart_system_architecture():
    """生成 系统架构与数据流图"""
    fig, ax = plt.subplots(1, 1, figsize=(12.8, 4.8))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 6)
    ax.axis('off')
    ax.set_facecolor('#FAFBFC')

    boxes = [
        # (x, y, w, h, text, color, text_color)
        (0.3, 2.3, 2.5, 1.4, '原始数据\n槽况数据_处理后_v2', MPL_LIGHT, MPL_NAVY),
        (3.3, 2.3, 2.5, 1.4, 'DataProcessor\n特征工程 + 标准化', '#D5E8D4', '#1B5E20'),
        (6.3, 3.7, 2.5, 1.4, '训练基础LSTM\nbest_model.pth', '#FCE4D6', '#7B3F00'),
        (6.3, 1.0, 2.5, 1.4, '训练条件预测器\nbest_conditional.pth', '#FCE4D6', '#7B3F00'),
        (9.3, 2.3, 2.5, 1.4, 'VoltageControlEnv\n仿真环境 (MDP)', '#E1D5E7', '#4A235A'),
        (12.3, 2.3, 2.5, 1.4, 'MPD-PPO 训练\nbest_ppo_model.pth', '#FFCCCC', '#7B0000'),
    ]

    for x, y, w, h, text, color, tc in boxes:
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.15",
                               facecolor=color, edgecolor=MPL_STEEL, linewidth=1.2)
        ax.add_patch(rect)
        # Split multi-line
        lines = text.split('\n')
        for j, line in enumerate(lines):
            fontsize = 10 if j == 0 else 8.5
            fontweight = 'bold' if j == 0 else 'normal'
            ypos = y + h - 0.3 - j * 0.35
            ax.text(x + w/2, ypos, line, ha='center', va='center',
                    fontsize=fontsize, fontweight=fontweight, color=tc)

    # Arrows between boxes
    arrows = [
        (2.8, 3.0, 3.3, 3.0),   # 原始数据 → DataProcessor
        (5.8, 3.0, 6.3, 4.4),   # DataProcessor → 基础LSTM
        (5.8, 3.0, 6.3, 1.7),   # DataProcessor → 条件预测器
        (8.8, 4.4, 9.3, 3.0),   # 基础LSTM → 条件预测器 (向下)
        (8.8, 1.7, 9.3, 3.0),   # 条件预测器 → Env
        (11.8, 3.0, 12.3, 3.0), # Env → PPO
    ]
    for x1, y1, x2, y2 in arrows:
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2.2,
                                    connectionstyle='arc3,rad=0'))

    # Title box at top
    ax.text(8, 5.6, '训练阶段 — 离线流水线', ha='center', va='center',
            fontsize=15, fontweight='bold', color=MPL_NAVY)

    # Inference phase at bottom
    infer_text = '推理阶段:  实时工况(7天) + 设定电压(14天) → Actor(PPO策略) → 28维动作轨迹 → Day-1动作执行 + 条件预测器验证'
    rect2 = FancyBboxPatch((0.3, 0.05), 15.4, 0.7, boxstyle="round,pad=0.1",
                            facecolor=MPL_NAVY, edgecolor=MPL_NAVY, linewidth=1)
    ax.add_patch(rect2)
    ax.text(8, 0.4, infer_text, ha='center', va='center', fontsize=9.5,
            color='white', fontweight='bold')

    plt.tight_layout(pad=0.2)
    path = os.path.join(CHARTS_DIR, 'architecture.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFBFC', edgecolor='none')
    plt.close(fig)
    return path


def chart_conditional_predictor():
    """生成 条件电压预测器架构图"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 4))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 5)
    ax.axis('off')
    ax.set_facecolor('#FAFBFC')

    # Input boxes
    inputs = [
        (0.5, 3.2, 3.0, 1.0, '过去7天特征 (B,7,feat_dim)\n+ pot_ids (B,)', MPL_LIGHT, MPL_NAVY),
        (0.5, 0.8, 3.0, 1.0, '未来14天动作 (B,14,2)\n[ALF加料量, 实际出铝量]', '#D5E8D4', '#1B5E20'),
    ]
    for x, y, w, h, text, color, tc in inputs:
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.12",
                               facecolor=color, edgecolor=MPL_STEEL, linewidth=1.2)
        ax.add_patch(rect)
        lines = text.split('\n')
        for j, line in enumerate(lines):
            ax.text(x + w/2, y + h - 0.25 - j * 0.3, line, ha='center', va='center',
                    fontsize=9, fontweight='bold' if j == 0 else 'normal', color=tc)

    # Processing blocks
    proc_blocks = [
        (4.2, 3.2, 2.8, 1.0, '冻结的 LSTM + 槽嵌入\n→ hidden (128维)', '#E1D5E7', '#4A235A'),
        (4.2, 0.8, 2.8, 1.0, 'Conv1D(k=3) → AvgPool\n→ FC (32维)', '#E1D5E7', '#4A235A'),
        (7.8, 2.0, 2.0, 1.6, 'Concat\n(160维)', '#FCE4D6', '#7B3F00'),
        (10.5, 2.0, 2.5, 1.6, 'FC(64→32→14)\n预测头', '#FFCCCC', '#7B0000'),
        (13.5, 2.0, 2.2, 1.6, '未来14天\n电压预测', MPL_LIGHT, MPL_GREEN),
    ]
    for x, y, w, h, text, color, tc in proc_blocks:
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.1",
                               facecolor=color, edgecolor=MPL_STEEL, linewidth=1.2)
        ax.add_patch(rect)
        lines = text.split('\n')
        for j, line in enumerate(lines):
            ax.text(x + w/2, y + h - 0.25 - j * 0.32, line, ha='center', va='center',
                    fontsize=9, fontweight='bold' if j == 0 else 'normal', color=tc)

    # Arrows
    ax.annotate('', xy=(4.2, 3.7), xytext=(3.5, 3.7), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2))
    ax.annotate('', xy=(4.2, 1.3), xytext=(3.5, 1.3), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2))
    ax.annotate('', xy=(7.8, 2.8), xytext=(7.0, 3.7), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2))
    ax.annotate('', xy=(7.8, 2.8), xytext=(7.0, 1.3), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2))
    ax.annotate('', xy=(10.5, 2.8), xytext=(9.8, 2.8), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2))
    ax.annotate('', xy=(13.5, 2.8), xytext=(13.0, 2.8), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2))

    ax.text(8, 4.6, 'ConditionalVoltagePredictor 架构（迁移学习设计）', ha='center', va='center',
            fontsize=14, fontweight='bold', color=MPL_NAVY)
    ax.text(8, 0.05, '关键设计：冻结预训练 LSTM + 新增轻量条件编码器 → 零破坏模型扩展', ha='center',
            va='center', fontsize=9, color=MPL_MID_GRAY, style='italic')

    plt.tight_layout(pad=0.2)
    path = os.path.join(CHARTS_DIR, 'conditional_predictor.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFBFC', edgecolor='none')
    plt.close(fig)
    return path


def chart_model_hierarchy():
    """生成 三模型层次关系图"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 3.6))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 4.5)
    ax.axis('off')
    ax.set_facecolor('#FAFBFC')

    levels = [
        (1.0, 2.8, 3.5, 1.0, 'LSTMModel (基础)\n输入: (B,7,feat_dim)\n输出: (B,14)', '#C8E6C9', '#1B5E20'),
        (6.2, 2.8, 4.0, 1.0, 'LSTMModelWithPotEmbedding (增强)\n输入: 基础 + pot_ids\n输出: (B,14) + BatchNorm', '#BBDEFB', '#0D47A1'),
        (11.0, 2.8, 4.2, 1.0, 'ConditionalVoltagePredictor (条件)\n输入: 增强 + 动作(B,14,2)\n输出: (B,14)', '#FFCDD2', '#B71C1C'),
    ]
    for x, y, w, h, text, color, tc in levels:
        rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.1",
                               facecolor=color, edgecolor=MPL_STEEL, linewidth=1.3)
        ax.add_patch(rect)
        lines = text.split('\n')
        for j, line in enumerate(lines):
            fs = 9.5 if j == 0 else 8
            fw = 'bold' if j == 0 else 'normal'
            ax.text(x + w/2, y + h - 0.2 - j * 0.28, line, ha='center', va='center',
                    fontsize=fs, fontweight=fw, color=tc)

    # Up arrows
    ax.annotate('', xy=(8.2, 2.8), xytext=(4.5, 3.8), arrowprops=dict(arrowstyle='->', color=MPL_GRAY, lw=2))
    ax.annotate('', xy=(13.0, 2.8), xytext=(10.2, 3.8), arrowprops=dict(arrowstyle='->', color=MPL_GRAY, lw=2))
    ax.text(5.8, 3.5, '继承', ha='center', fontsize=8, color=MPL_GRAY)
    ax.text(11.8, 3.5, '继承', ha='center', fontsize=8, color=MPL_GRAY)

    # Label
    ax.text(8, 4.2, '模型层次：基础 → 增强 → 条件（逐层继承扩展）', ha='center', fontsize=13,
            fontweight='bold', color=MPL_NAVY)

    # Training note
    ax.text(8, 0.3, '迁移学习策略：基础 LSTM 加载预训练权重后冻结 → 仅训练条件编码器和预测头',
            ha='center', fontsize=9, color=MPL_MID_GRAY, style='italic')

    plt.tight_layout(pad=0.2)
    path = os.path.join(CHARTS_DIR, 'model_hierarchy.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFBFC', edgecolor='none')
    plt.close(fig)
    return path


def chart_ppo_algorithm():
    """生成 MPD-PPO 算法流程图"""
    fig, ax = plt.subplots(1, 1, figsize=(12.8, 4.8))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 6)
    ax.axis('off')
    ax.set_facecolor('#FAFBFC')

    # State input
    state_box = FancyBboxPatch((0.3, 3.4), 2.5, 1.2, boxstyle="round,pad=0.15",
                                facecolor=MPL_LIGHT, edgecolor=MPL_STEEL, linewidth=1.2)
    ax.add_patch(state_box)
    ax.text(1.55, 4.4, '状态 s(t)', ha='center', fontsize=10, fontweight='bold', color=MPL_NAVY)
    ax.text(1.55, 3.9, '7天特征 + 设定电压\n+ 槽号', ha='center', fontsize=8, color=MPL_NAVY)

    # Actor
    actor_box = FancyBboxPatch((3.5, 3.4), 2.5, 1.2, boxstyle="round,pad=0.15",
                                facecolor='#D5E8D4', edgecolor=MPL_GREEN, linewidth=1.2)
    ax.add_patch(actor_box)
    ax.text(4.75, 4.4, 'Actor 网络', ha='center', fontsize=10, fontweight='bold', color='#1B5E20')
    ax.text(4.75, 3.9, 'LSTM编码 → μ(28d)\nσ(28d) → 采样', ha='center', fontsize=8, color='#1B5E20')

    # Action
    action_box = FancyBboxPatch((6.8, 3.4), 2.5, 1.2, boxstyle="round,pad=0.15",
                                 facecolor='#FCE4D6', edgecolor=MPL_ORANGE, linewidth=1.2)
    ax.add_patch(action_box)
    ax.text(8.05, 4.4, '28维动作轨迹', ha='center', fontsize=10, fontweight='bold', color='#7B3F00')
    ax.text(8.05, 3.9, '[alf1, out1, ..., alf14, out14]\n反归一化 + 裁剪 + 约束', ha='center', fontsize=8, color='#7B3F00')

    # Environment
    env_box = FancyBboxPatch((10.1, 3.4), 2.5, 1.2, boxstyle="round,pad=0.15",
                              facecolor='#E1D5E7', edgecolor='#4A235A', linewidth=1.2)
    ax.add_patch(env_box)
    ax.text(11.35, 4.4, 'VoltageControlEnv', ha='center', fontsize=10, fontweight='bold', color='#4A235A')
    ax.text(11.35, 3.9, '预测 → 奖励 → 滑动\n窗口 → s(t+1)', ha='center', fontsize=8, color='#4A235A')

    # Reward
    reward_box = FancyBboxPatch((13.4, 3.4), 2.3, 1.2, boxstyle="round,pad=0.15",
                                 facecolor='#FFCCCC', edgecolor=MPL_RED, linewidth=1.2)
    ax.add_patch(reward_box)
    ax.text(14.55, 4.4, '多步加权奖励', ha='center', fontsize=10, fontweight='bold', color='#7B0000')
    ax.text(14.55, 3.9, 'R = R_acc + P_smooth\n+ P_bound', ha='center', fontsize=8, color='#7B0000')

    # Arrows
    for (x1, x2, y) in [(2.8, 3.5, 4.0), (6.0, 6.8, 4.0), (9.3, 10.1, 4.0), (12.6, 13.4, 4.0)]:
        ax.annotate('', xy=(x2, y), xytext=(x1, y), arrowprops=dict(arrowstyle='->', color=MPL_STEEL, lw=2.2))

    # Critic path (bottom)
    critic_box = FancyBboxPatch((3.5, 1.6), 2.5, 1.0, boxstyle="round,pad=0.15",
                                 facecolor='#BBDEFB', edgecolor=MPL_STEEL, linewidth=1.2)
    ax.add_patch(critic_box)
    ax.text(4.75, 2.3, 'Critic 网络', ha='center', fontsize=10, fontweight='bold', color='#0D47A1')
    ax.text(4.75, 1.8, 'LSTM编码 → V(s)∈[-50,50]', ha='center', fontsize=8, color='#0D47A1')

    # GAE
    gae_box = FancyBboxPatch((6.8, 1.6), 2.5, 1.0, boxstyle="round,pad=0.15",
                              facecolor='#BBDEFB', edgecolor=MPL_STEEL, linewidth=1.2)
    ax.add_patch(gae_box)
    ax.text(8.05, 2.3, 'GAE 优势估计', ha='center', fontsize=10, fontweight='bold', color='#0D47A1')
    ax.text(8.05, 1.8, 'A = Σ (γλ)^t · δ(t+k)\n标准化 + TD回退', ha='center', fontsize=8, color='#0D47A1')

    # PPO Update
    ppo_box = FancyBboxPatch((10.1, 1.6), 2.5, 1.0, boxstyle="round,pad=0.15",
                              facecolor='#BBDEFB', edgecolor=MPL_STEEL, linewidth=1.2)
    ax.add_patch(ppo_box)
    ax.text(11.35, 2.3, 'PPO 更新(×10轮)', ha='center', fontsize=10, fontweight='bold', color='#0D47A1')
    ax.text(11.35, 1.8, 'ALF(ε=0.1) + OUT(ε=0.2)\n差异化裁剪+求和', ha='center', fontsize=8, color='#0D47A1')

    # Arrows for critic path
    ax.annotate('', xy=(3.5, 2.1), xytext=(1.55, 2.1), arrowprops=dict(arrowstyle='->', color='#607D8B', lw=2))
    ax.annotate('', xy=(6.8, 2.1), xytext=(6.0, 2.1), arrowprops=dict(arrowstyle='->', color='#607D8B', lw=2))
    ax.annotate('', xy=(10.1, 2.1), xytext=(9.3, 2.1), arrowprops=dict(arrowstyle='->', color='#607D8B', lw=2))

    # Flow loop (update actor)
    ax.annotate('', xy=(11.35, 2.6), xytext=(4.75, 3.4),
                arrowprops=dict(arrowstyle='->', color=MPL_RED, lw=1.8,
                                connectionstyle='arc3,rad=-0.3', ls='--'))
    ax.text(8, 2.8, '更新 Actor/Critic', ha='center', fontsize=8, color=MPL_RED, style='italic')

    # Title
    ax.text(8, 5.5, 'MPD-PPO 强化学习算法 — 训练循环', ha='center', fontsize=14, fontweight='bold', color=MPL_NAVY)

    # Bottom note
    ax.text(8, 0.3, '收集256步经验 → GAE计算优势 → 10轮mini-batch(64) PPO更新 → 清空Buffer → 循环',
            ha='center', fontsize=9, color=MPL_MID_GRAY, style='italic')

    plt.tight_layout(pad=0.2)
    path = os.path.join(CHARTS_DIR, 'ppo_algorithm.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFBFC', edgecolor='none')
    plt.close(fig)
    return path


def chart_improvement_timeline():
    """生成 改进方案时间线图"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 3.2))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 4)
    ax.axis('off')
    ax.set_facecolor('#FAFBFC')

    # Timeline axis
    ax.plot([1, 15], [2, 2], '-', color=MPL_STEEL, lw=3, zorder=1)

    # R1
    ax.plot(4, 2, 'o', color=MPL_GREEN, markersize=14, zorder=2)
    ax.text(4, 2.8, '第一轮重构\n"滚动时域轨迹 PPO"', ha='center', fontsize=10, fontweight='bold', color=MPL_GREEN)
    ax.text(4, 2.6, '修改 1-6：动作空间 2维→28维\n滚动时域 + 多步奖励 + 差异化裁剪', ha='center', fontsize=7.5, color=MPL_DARK_GRAY)
    ax.plot([4, 4], [2, 2.55], '-', color=MPL_GREEN, lw=1.5, zorder=1)

    # R2
    ax.plot(11, 2, 'o', color=MPL_RED, markersize=14, zorder=2)
    ax.text(11, 2.8, '第二轮修复\n"训练稳定性 Bug 修复"', ha='center', fontsize=10, fontweight='bold', color=MPL_RED)
    ax.text(11, 2.6, '修复 7-12：形状崩溃 / 惩罚爆炸\nGAE坍缩 / 值函数漂移 / 初始探索', ha='center', fontsize=7.5, color=MPL_DARK_GRAY)
    ax.plot([11, 11], [2, 2.55], '-', color=MPL_RED, lw=1.5, zorder=1)

    # Labels
    ax.text(4, 1.2, '架构升级', ha='center', fontsize=9, color=MPL_MID_GRAY, style='italic')
    ax.text(11, 1.2, '稳定性修复', ha='center', fontsize=9, color=MPL_MID_GRAY, style='italic')

    # Before/After
    ax.text(1.5, 3.5, '原始方案', ha='center', fontsize=10, fontweight='bold', color=MPL_GRAY)
    ax.text(1.5, 3.2, '单步交互 PPO\n2维动作空间', ha='center', fontsize=7.5, color=MPL_GRAY)

    ax.text(14.5, 3.5, '当前状态', ha='center', fontsize=10, fontweight='bold', color=MPL_STEEL)
    ax.text(14.5, 3.2, '28维轨迹 PPO\n训练稳定', ha='center', fontsize=7.5, color=MPL_STEEL)

    plt.tight_layout(pad=0.2)
    path = os.path.join(CHARTS_DIR, 'improvement_timeline.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFBFC', edgecolor='none')
    plt.close(fig)
    return path


def chart_reward_composition():
    """生成 奖励函数组成结构图"""
    fig, ax = plt.subplots(1, 1, figsize=(12, 3.2))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 4)
    ax.axis('off')
    ax.set_facecolor('#FAFBFC')

    # Three components
    components = [
        (1.0, 1.0, 4.2, 'R_acc — 精度奖励\n10.0 × Σ w_i × exp(-2.0 × (v_pred_i - v_set_i)²)\nw_i = [1.0, 0.9, ..., 0.1]\n指数衰减 + 时间衰减权重', MPL_GREEN, '#1B5E20'),
        (5.9, 1.0, 4.2, 'P_smooth — 平滑违规惩罚\n-5.0 × Σ max(0, |Δalf|/15 - 0.2)\n-5.0 × Σ max(0, |Δout|/700 - 0.143)\n归一化到 [0,1] 比例空间', MPL_ORANGE, '#7B3F00'),
        (10.8, 1.0, 4.2, 'P_bound — 边界软约束\n-10 × 超出比例\nALF∈[25,40] kg\nOUT∈[3500,4200] kg\n软约束避免硬截断', MPL_RED, '#7B0000'),
    ]
    for x, y, w, text, c1, c2 in components:
        rect = FancyBboxPatch((x, y), w, 2.2, boxstyle="round,pad=0.12",
                               facecolor=c1, edgecolor=c2, linewidth=1.3, alpha=0.15)
        ax.add_patch(rect)
        lines = text.split('\n')
        for j, line in enumerate(lines):
            fs = 10 if j == 0 else 8
            fw = 'bold' if j == 0 else 'normal'
            ax.text(x + w/2, y + 2.0 - j * 0.3, line, ha='center', va='center',
                    fontsize=fs, fontweight=fw, color=c2)

    # Plus signs
    ax.text(5.55, 2.1, '+', ha='center', fontsize=18, fontweight='bold', color=MPL_GRAY)
    ax.text(10.45, 2.1, '+', ha='center', fontsize=18, fontweight='bold', color=MPL_GRAY)

    # Total
    ax.text(8, 3.6, 'R = R_acc + P_smooth + P_bound', ha='center', fontsize=14, fontweight='bold', color=MPL_NAVY)
    ax.text(8, 3.3, '单步奖励合理范围 ≈ [-30, +70]（修复后）', ha='center', fontsize=9, color=MPL_MID_GRAY, style='italic')

    plt.tight_layout(pad=0.2)
    path = os.path.join(CHARTS_DIR, 'reward_composition.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='#FAFBFC', edgecolor='none')
    plt.close(fig)
    return path


# ── PPT 构建函数 ────────────────────────────────────────

def add_bg_rect(slide, left, top, width, height, color):
    """添加背景矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_navy_header(slide, title_text, slide_num=None):
    """添加统一的深海军蓝页眉"""
    # Top bar
    bar = add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), NAVY)
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Microsoft YaHei'
    # Page number
    if slide_num:
        numBox = slide.shapes.add_textbox(Inches(12), Inches(0.2), Inches(1), Inches(0.6))
        nf = numBox.text_frame
        np = nf.paragraphs[0]
        np.text = f'{slide_num:02d}'
        np.font.size = Pt(18)
        np.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
        np.font.name = 'Calibri'
        np.alignment = PP_ALIGN.RIGHT
    # Bottom accent line
    line = add_bg_rect(slide, Inches(0.6), Inches(1.05), Inches(2.5), Pt(4), ACCENT_R)

def add_body_text(slide, left, top, width, height, paragraphs_data):
    """添加正文文本区域
    paragraphs_data: list of (text, font_size, bold, color, alignment)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, fs, bold, color, align) in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = fs  # fs is already a Pt object from caller
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_card(slide, left, top, width, height, title, lines, title_color=None, bg_color=None):
    """添加卡片式内容块"""
    if bg_color is None:
        bg_color = LIGHT_BG
    if title_color is None:
        title_color = STEEL

    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = LIGHT_LINE
    shape.line.width = Pt(1)

    # Title
    txBox = slide.shapes.add_textbox(Inches(left.inches + 0.15), Inches(top.inches + 0.1),
                                      Inches(width.inches - 0.3), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = 'Microsoft YaHei'

    # Content lines
    txBox2 = slide.shapes.add_textbox(Inches(left.inches + 0.15), Inches(top.inches + 0.45),
                                       Inches(width.inches - 0.3), Inches(height.inches - 0.55))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(9.5)
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(3)

def add_table(slide, left, top, col_widths, headers, rows, header_bg=None):
    """添加格式化表格"""
    if header_bg is None:
        header_bg = NAVY
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, Inches(total_w), Inches(0.35 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg

    # Data rows
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else LIGHT_BG
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8.5)
                p.font.color.rgb = DARK_GRAY
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

def add_image_slide(slide, img_path, left, top, width, height=None):
    """插入图片"""
    if height is None:
        return slide.shapes.add_picture(img_path, left, top, width)
    return slide.shapes.add_picture(img_path, left, top, width, height)

def add_bottom_note(slide, text, color=None):
    """添加页脚注释"""
    if color is None:
        color = MID_GRAY
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'


# ── 构建 PPT ────────────────────────────────────────────

def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ═══════════════════════════════════════════════════════
    # Slide 1 — 封面
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    # Full navy background
    add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    # Accent bar
    add_bg_rect(slide, Inches(0.6), Inches(3.15), Inches(0.08), Inches(1.6), ACCENT_R)

    add_body_text(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(1.5), [
        ('铝电解槽工作电压参数优化', Pt(40), True, WHITE, PP_ALIGN.LEFT),
        ('项目代码分析总览', Pt(28), False, RGBColor(0xAA, 0xBB, 0xDD), PP_ALIGN.LEFT),
    ])
    add_body_text(slide, Inches(1.0), Inches(3.3), Inches(11), Inches(1.0), [
        ('条件电压预测 + PPO 强化学习优化  |  滚动时域轨迹 PPO  |  MPD-PPO 算法', Pt(14), False, RGBColor(0x88, 0x99, 0xBB), PP_ALIGN.LEFT),
    ])
    # Bottom info
    add_body_text(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.8), [
        ('基于 Python / PyTorch / LSTM / PPO 的工业参数优化系统', Pt(12), False, RGBColor(0x77, 0x88, 0xAA), PP_ALIGN.LEFT),
        ('2026年5月', Pt(11), False, RGBColor(0x66, 0x77, 0x99), PP_ALIGN.LEFT),
    ])

    # ═══════════════════════════════════════════════════════
    # Slide 2 — 目录
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '目  录', 2)

    toc_items = [
        ('01', '项目概览', '核心目标、技术路线与项目定位'),
        ('02', '目录结构与角色总览', '项目文件树、模块职责划分'),
        ('03', '模块A：预测模型 — 配置与数据', 'config.py 超参数、6步数据处理流水线'),
        ('04', '模块A：预测模型 — 网络架构', '三模型层次、条件预测器、迁移学习'),
        ('05', '模块B：PPO — 配置与环境', 'PPO 超参数、MDP 仿真环境、奖励函数'),
        ('06', '模块B：PPO — 算法实现', 'Actor/Critic 网络、MPD-PPO、差异化裁剪'),
        ('07', '系统架构与数据流', '训练/推理整体流水线、模型依赖关系'),
        ('08', '第一轮改进：架构重构', '修改 1-6：28维轨迹、多步奖励、滚动时域'),
        ('09', '第二轮改进：稳定性修复', '修复 7-12：奖励归一化、GAE 防坍缩'),
        ('10', '项目优势总结', '5大核心优势与技术亮点'),
        ('11', '潜在改进方向', '5个未来扩展方向'),
        ('12', '总结与展望', '全文总结与致谢'),
    ]
    for i, (num, title, desc) in enumerate(toc_items):
        row = i // 2
        col = i % 2
        lx = Inches(0.6 + col * 6.2)
        ty = Inches(1.5 + row * 0.85)
        # Number circle
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx, ty, Inches(0.45), Inches(0.45))
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY if i < 2 else STEEL
        circ.line.fill.background()
        tf = circ.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        # Title + desc
        add_body_text(slide, Inches(lx.inches + 0.6), Inches(ty.inches + 0.05),
                       Inches(5.2), Inches(0.7), [
            (title, Pt(14), True, NAVY, PP_ALIGN.LEFT),
            (desc, Pt(9), False, MID_GRAY, PP_ALIGN.LEFT),
        ])

    # ═══════════════════════════════════════════════════════
    # Slide 3 — 项目概览
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '项目概览', 3)

    add_card(slide, Inches(0.6), Inches(1.4), Inches(6.0), Inches(2.6),
             '核心目标', [
                 '• 智能调整 ALF加料量 和 实际出铝量 两个可控参数',
                 '• 使电解槽工作平均电压尽可能逼近设定电压',
                 '• 提升电流效率、降低能耗',
                 '',
                 '技术路线：条件电压预测 + PPO强化学习优化',
                 '核心问题：「给定过去7天工况 + 未来14天操作计划',
                 '                 → 未来14天电压会是多少？」',
             ], ACCENT_R)

    add_card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.6),
             '项目演进', [
                 '• 原始方案：单步交互 PPO，2维动作空间',
                 '• 第一轮重构：升级为「滚动时域轨迹 PPO」',
                 '    动作空间 2维 → 28维（14天×2动作）',
                 '• 第二轮修复：训练稳定性 Bug 修复',
                 '    （奖励归一化、GAE 防坍缩、值函数约束）',
                 '',
                 '当前状态：✅ 12项改进全部实施完毕',
             ], ACCENT_G)

    add_card(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5),
             '关键技术指标', [
                 '• 预测模型：2层 LSTM (hidden=128) + 槽号16维嵌入，输入7天→预测14天，MAE/RMSE/MAPE/R² 多指标评估',
                 '• PPO 训练：每轮收集256步经验，10轮内层迭代×mini_batch(64)，差异化裁剪 ALF(ε=0.1) / OUT(ε=0.2)',
                 '• 奖励函数：R = R_acc(10.0×指数精度) + P_smooth(-5.0×归一化违规) + P_bound(-10×边界软约束)，单步约[-30,+70]',
                 '• 数据划分：训练30槽/验证6槽/测试6槽，按槽号独立划分避免数据泄露；特征维度约84维（12基础+60统计+4衍生+8目标）',
             ], STEEL)

    # ═══════════════════════════════════════════════════════
    # Slide 4 — 目录结构与角色
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '目录结构与角色总览', 4)

    # Left side: directory tree
    dir_tree = [
        'model/                          模块A：条件电压预测模型',
        '  ├── config.py                  预测模型超参数 + 特征配置',
        '  ├── data_processor.py          数据加载/特征工程/序列构建',
        '  ├── model.py                   LSTM + 条件预测器定义',
        '  ├── train.py                   预测模型训练脚本',
        '  ├── predict.py                 预测模型推理脚本',
        '  └── output/                    训练产出（权重、图表、scaler）',
        '',
        'ppo参数优化/model/               模块B：PPO强化学习优化',
        '  ├── config.py                  PPO超参数 + 动作/奖励配置',
        '  ├── environment.py             电压控制仿真环境 (MDP)',
        '  ├── ppo.py                     Actor/Critic + MPD-PPO算法',
        '  ├── train_ppo.py               PPO训练脚本',
        '  └── visualize_ppo.py           PPO结果可视化',
    ]
    add_body_text(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.8),
                  [(line, Pt(8.5), False, DARK_GRAY, PP_ALIGN.LEFT) for line in dir_tree])

    # Right side: role descriptions
    add_card(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(1.8),
             '模块A：基础层（预测器）', [
                 '回答核心问题：给定过去7天工况 + 未来14天操作计划',
                 '→ 未来14天电压会是多少？',
                 'LSTM时序模型 + 槽号嵌入 + 迁移学习',
                 '为PPO提供仿真环境的基础能力',
             ], STEEL)

    add_card(slide, Inches(6.9), Inches(3.5), Inches(5.8), Inches(1.8),
             '模块B：决策层（优化器）', [
                 '在预测模型提供的仿真环境中训练PPO智能体',
                 'Actor输出28维动作轨迹（14天×2动作）',
                 '差异化裁剪 + GAE优势估计 + 多步加权奖励',
                 '使智能体学会生成最优的14天操作计划',
             ], ACCENT_R)

    add_card(slide, Inches(6.9), Inches(5.6), Inches(5.8), Inches(0.9),
             '参考文档 & 分析总结', [
                 'agents/：方案设计文档（架构演进记录）',
                 'result/：分析总结 + 改进方案跟踪',
             ], ACCENT_O)

    # ═══════════════════════════════════════════════════════
    # Slide 5 — 模块A：配置与数据
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '模块A：预测模型 — 配置与数据处理', 5)

    add_table(slide, Inches(0.6), Inches(1.4), [2.0, 1.5, 2.5],
              ['参数类别', '关键参数', '说明'],
              [
                  ['序列配置', 'INPUT_LEN=7, OUTPUT_LEN=14', '输入7天，预测14天'],
                  ['模型结构', 'HIDDEN_DIM=128, NUM_LAYERS=2', '2层LSTM + 槽号16维嵌入'],
                  ['训练控制', 'BATCH_SIZE=64, LR=0.001', 'Adam优化器，ReduceLROnPlateau'],
                  ['早停策略', 'PATIENCE=15', '验证集loss不降15轮即停止'],
                  ['数据划分', '30训练/6验证/6测试槽', '按槽号独立划分，避免数据泄露'],
                  ['特征选择', '12个高相关性特征(|r|≥0.2)', '基于Spearman相关系数筛选'],
              ])

    add_table(slide, Inches(7.5), Inches(1.4), [2.0, 3.5],
              ['处理步骤', '说明'],
              [
                  ['1. load_data()', '读取Excel，按日期排序'],
                  ['2. preprocess_data()', '缺失值：分组前向填充+线性插值+均值填充'],
                  ['3. create_features()', '特征工程：滑动统计(3/7天)+差分(1/7阶)+衍生特征'],
                  ['4. 数据划分', 'time模式(6:2:2) / pot模式(按槽号分组)'],
                  ['5. 标准化', '仅在训练集fit scaler，然后transform全部'],
                  ['6. 序列构建', '滑动窗口切分(7天输入+14天输出)，最终≈84维特征'],
              ])

    add_card(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2.3),
             '特征工程详情（关键设计）', [
                 '• 基础特征(12)：基于 Spearman |r|≥0.2 筛选的高相关性工况参数',
                 '• 滑动统计(60)：每个基础特征 × 5种统计（3天mean/std、7天mean/std、完整统计），捕获近期趋势和波动',
                 '• 差分特征：1阶差分(日变化速度) + 7阶差分(周变化趋势)',
                 '• 衍生特征(4)：电压偏差、铝电解比例、槽龄_log、槽龄_squared — 引入领域知识',
                 '• 目标变量统计(8)：工作平均电压等目标变量的3/7天均值/标准差',
                 '• 条件预测专用：use_future_actions=True 时额外提取未来14天[ALF加料量, 实际出铝量]序列 (B,14,2)',
             ], STEEL)

    # ═══════════════════════════════════════════════════════
    # Slide 6 — 模块A：网络架构
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '模块A：预测模型 — 网络架构', 6)

    # Model hierarchy chart
    chart_path = chart_model_hierarchy()
    add_image_slide(slide, chart_path, Inches(0.5), Inches(1.35), Inches(12.3))

    # Conditional predictor detail
    chart2_path = chart_conditional_predictor()
    add_image_slide(slide, chart2_path, Inches(0.3), Inches(4.35), Inches(12.7))

    add_bottom_note(slide, '关键设计：冻结预训练 LSTM + 条件编码器(Conv1D→AvgPool→FC) + 轻量预测头(FC 64→32→14) = 零破坏模型扩展')

    # ═══════════════════════════════════════════════════════
    # Slide 7 — 模块B：PPO 配置与环境
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '模块B：PPO — 配置与环境', 7)

    add_table(slide, Inches(0.6), Inches(1.4), [2.2, 2.0, 2.2],
              ['类别', '参数', '说明'],
              [
                  ['动作空间', 'ACTION_TRAJECTORY_DIM=28', '14天×2动作完整轨迹（核心改进）'],
                  ['PPO训练', 'STEPS_PER_UPDATE=256', '每次更新前收集256步经验'],
                  ['', 'MINI_BATCH_SIZE=64', '内层mini-batch大小'],
                  ['', 'INNER_EPOCHS=10', '每次更新10轮内层迭代'],
                  ['', 'LR=3e-4, γ=0.99, λ=0.95', '学习率/折扣/GAE参数'],
                  ['差异化裁剪', 'EPS_CLIP_ALF=0.1', 'ALF敏感，小步调整'],
                  ['', 'EPS_CLIP_OUT=0.2', '出铝量可稍大步'],
                  ['动作约束', 'ALF∈[25,40], OUT∈[3500,4200]', 'kg单位，含日变化上限'],
                  ['经验系数', 'EMPIRICAL_AL_OUT_RATIO=-0.3', '铝水平经验更新系数'],
              ])

    # Reward function
    add_card(slide, Inches(7.5), Inches(1.4), Inches(5.2), Inches(2.6),
             '奖励函数（已修复）', [
                 'R = R_acc + P_smooth + P_bound',
                 '',
                 'R_acc = 10.0 × Σ w_i × exp(-2.0 × (v_pred_i - v_set_i)²)',
                 '    w_i = [1.0, 0.9, ..., 0.1]  (14天时间衰减)',
                 '',
                 'P_smooth = -5.0 × Σ violation_ratio',
                 '    (|Δalf|/15 或 |Δout|/700，归一化到[0,1])',
                 '',
                 'P_bound = -10 × 超出比例',
                 '单步奖励 ≈ [-30, +70]（修复后合理范围）',
             ], ACCENT_R)

    # Env details
    add_card(slide, Inches(0.6), Inches(4.7), Inches(6.6), Inches(2.2),
             'MDP 状态空间', [
                 'state_dim = 7 × input_dim + 14 + 1',
                 '• 过去7天完整特征窗口（展平）',
                 '• 未来14天目标设定电压序列',
                 '• 槽号索引',
                 '',
                 '状态转移（滚动时域）：',
                 '28维轨迹 → 重塑(14,2) → 预测14天电压',
                 '→ 仅执行day-1动作 → 滑动窗口1天',
                 '→ 经验更新非控特征 → 重算统计特征',
             ], STEEL)

    add_card(slide, Inches(7.5), Inches(4.7), Inches(5.2), Inches(2.2),
             '经验更新规则', [
                 '• 工作平均电压：用预测值更新',
                 '• ALF/出铝量：用 day-1 动作更新',
                 '• 铝水平：经验公式更新',
                 '   -0.3×(Δout/prev_out)×aluminum_level',
                 '• 电压设定/实际设定：沿用前一天',
                 '• 统计特征：基于新窗口重新计算',
                 '   3/7天均值/标准差/差分',
             ], ACCENT_O)

    # ═══════════════════════════════════════════════════════
    # Slide 8 — PPO 算法实现
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '模块B：PPO — 算法实现', 8)

    # PPO algorithm chart
    chart_path = chart_ppo_algorithm()
    add_image_slide(slide, chart_path, Inches(0.3), Inches(1.35), Inches(12.7))

    # Key details
    add_card(slide, Inches(0.6), Inches(5.15), Inches(5.8), Inches(1.8),
             'Actor / Critic 网络结构', [
                 'Actor: 状态 → LSTM+槽嵌入 → SharedNet → mean_head(28d)+std_head(28d)',
                 '  28个独立高斯分布，tanh缩放输出到[-1,1]，std_head bias=-1.0',
                 '  内置旧版兼容映射（4头→2头参数自动迁移）',
                 'Critic: 与Actor共享LSTM编码器结构 → 输出标量V(s)∈[-50,50]',
                 '  从[-1000,1000]收紧，防止值函数漂移导致GAE崩塌',
             ], STEEL)

    add_card(slide, Inches(6.9), Inches(5.15), Inches(5.8), Inches(1.8),
             '差异化裁剪（核心算法改进）', [
                 '独立裁剪 + 求和（而非全局裁剪）：',
                 '  ALF(偶数索引): ratio × A vs clip(ratio, 0.9, 1.1) × A → min',
                 '  OUT(奇数索引): ratio × A vs clip(ratio, 0.8, 1.2) × A → min',
                 '  actor_loss = -(alf_surr + out_surr)',
                 '',
                 'PPO更新流程：rewards标准化 → GAE优势 → 10轮mini-batch',
                 '  独立裁剪ALF/OUT → MSE计算critic_loss → 梯度裁剪1.0',
                 '  更新后自动清空经验缓冲区',
             ], ACCENT_R)

    # ═══════════════════════════════════════════════════════
    # Slide 9 — 系统架构与数据流
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '系统架构与数据流', 9)

    chart_path = chart_system_architecture()
    add_image_slide(slide, chart_path, Inches(0.2), Inches(1.3), Inches(12.9))

    # Model dependency
    add_card(slide, Inches(0.6), Inches(5.2), Inches(6.0), Inches(1.8),
             '模型间依赖关系', [
                 'LSTMModelWithPotEmbedding (基础时序编码器)',
                 '  ├→ ConditionalVoltagePredictor (条件预测=基础编码器+条件编码器)',
                 '  │    └→ VoltageControlEnv (仿真环境，封装预测器为MDP)',
                 '  │         └→ MPDPPO (含Actor/Critic，在仿真环境中训练)',
                 '  └── Actor/Critic 内部也有独立LSTM编码器（结构对称，参数独立）',
             ], STEEL)

    add_card(slide, Inches(6.9), Inches(5.2), Inches(5.8), Inches(1.8),
             '推理部署流程', [
                 '1. 实时工况(过去7天) + 设定电压(未来14天)',
                 '2. Actor(PPO策略) → 28维动作轨迹',
                 '3. Day-1动作 → 现场执行/人工审核',
                 '4. 条件预测器验证 → 预期未来14天电压曲线',
                 '',
                 '关键设计：训练-部署一致性',
                 '  Actor输出完整14天轨迹 → 预测器输入在训练',
                 '  和推理时保持一致 → 消除分布偏移',
             ], ACCENT_G)

    # ═══════════════════════════════════════════════════════
    # Slide 10 — 第一轮改进
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '第一轮改进：PPO 架构重构（修改 1-6）', 10)

    # Timeline chart
    chart_path = chart_improvement_timeline()
    add_image_slide(slide, chart_path, Inches(0.4), Inches(1.35), Inches(12.5))

    # Details
    mod1_data = [
        ('修改1', '优化目标切换', '从历史工作电压切换到设定电压', 'train_ppo.py:136', ACCENT_G),
        ('修改2', '28维轨迹输出', 'Actor输出14天×2动作 + 独立裁剪求和', 'ppo.py, config.py', ACCENT_G),
        ('修改3', '多步加权奖励', '14天时间衰减指数奖励函数', 'environment.py:_calculate_reward()', ACCENT_G),
        ('修改4', '滚动时域交互', '环境step执行完整轨迹，仅推进1天', 'environment.py:step()', ACCENT_G),
        ('修改5', '非控特征更新', '铝水平/统计特征的经验更新规则', 'environment.py:_update_state()', ACCENT_G),
        ('修改6', 'PPO更新逻辑', '差异化裁剪+mini-batch+GAE边界', 'ppo.py:update()', ACCENT_G),
    ]
    for i, (mod_id, name, desc, loc, clr) in enumerate(mod1_data):
        col = i % 3
        row = i // 3
        lx = Inches(0.6 + col * 4.15)
        ty = Inches(3.55 + row * 1.85)
        add_card(slide, lx, ty, Inches(3.9), Inches(1.65), f'{mod_id}：{name}', [
            f'内容：{desc}',
            f'涉及：{loc}',
        ], clr, LIGHT_BG)

    # ═══════════════════════════════════════════════════════
    # Slide 11 — 第二轮改进 + 奖励函数
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_navy_header(slide, '第二轮改进：训练稳定性 Bug 修复（修复 7-12）', 11)

    # Reward composition chart
    chart_path = chart_reward_composition()
    add_image_slide(slide, chart_path, Inches(0.3), Inches(1.35), Inches(12.7))

    # Bug fixes in cards
    fixes_data = [
        ('修复7', 'target_voltage 形状崩溃', '滑动窗口用末值填充保持14维', ACCENT_R),
        ('修复8', '平滑惩罚量级爆炸', 'OUT惩罚-2260/step → 归一化到[0,1]', ACCENT_R),
        ('修复9', 'GAE 硬裁剪信号消失', '移除[-50,50]硬裁剪 → reward标准化', ACCENT_R),
        ('修复10', '值函数漂移(+1000)', 'Critic输出[-1000,1000]→[-50,50] + TD回退', ACCENT_R),
        ('修复11', '初始探索过度随机', 'std_head bias初始-1.0（std≈0.31而非0.694）', ACCENT_R),
        ('修复12', '缺乏诊断手段', 'info增加R_acc/P_smooth/P_bound分量', ACCENT_R),
    ]
    for i, (fid, problem, fix, clr) in enumerate(fixes_data):
        col = i % 3
        row = i // 3
        lx = Inches(0.6 + col * 4.15)
        ty = Inches(3.7 + row * 1.75)
        add_card(slide, lx, ty, Inches(3.9), Inches(1.5), f'{fid}：{problem}', [
            f'修复：{fix}',
        ], clr, LIGHT_BG)

    add_bottom_note(slide, '根因链条：平滑惩罚(OUT)量级放大 → reward全部极负 → hard clip抹除差异 → GAE崩塌 → Actor Loss=0 → 训练死锁。五项修复形成闭环。')

    # ═══════════════════════════════════════════════════════
    # Slide 12 — 总结与展望
    # ═══════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    # Navy background for closing
    add_bg_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_bg_rect(slide, Inches(0.6), Inches(3.15), Inches(0.08), Inches(1.6), ACCENT_R)

    add_body_text(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(1.5), [
        ('谢谢', Pt(42), True, WHITE, PP_ALIGN.LEFT),
        ('铝电解槽工作电压参数优化项目 — 代码分析总览', Pt(22), False, RGBColor(0xAA, 0xBB, 0xDD), PP_ALIGN.LEFT),
    ])

    # Core summary points
    summary_lines = [
        '❶  构建了完整「条件电压预测 + PPO强化学习优化」流水线',
        '❷  条件预测器采用迁移学习设计，冻结预训练LSTM + 轻量条件编码器',
        '❸  PPO智能体输出28维14天轨迹，差异化裁剪 ALF(ε=0.1)/OUT(ε=0.2)',
        '❹  多步加权奖励让Agent考虑长期电压质量，非贪心单日优化',
        '❺  两轮12项改进：架构重构（6项）+ 稳定性修复（6项），全部实施完毕',
    ]
    add_body_text(slide, Inches(1.0), Inches(3.4), Inches(11), Inches(2.0),
                  [(line, Pt(11), False, RGBColor(0xCC, 0xDD, 0xEE), PP_ALIGN.LEFT) for line in summary_lines])

    # Future directions
    add_body_text(slide, Inches(1.0), Inches(5.3), Inches(11), Inches(0.8), [
        ('改进方向：预测器动作编码升级 (Conv1D→Transformer) | 经验系数数据标定 | 多目标优化 | 在线学习 | Actor/Critic权重共享', Pt(9), False, RGBColor(0x88, 0x99, 0xBB), PP_ALIGN.LEFT),
    ])
    add_body_text(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.5), [
        ('Python / PyTorch / LSTM / PPO  |  2026年5月', Pt(10), False, RGBColor(0x77, 0x88, 0xAA), PP_ALIGN.LEFT),
    ])

    # ── 保存 ─────────────────────────────────────────────
    output_path = os.path.join(OUTPUT_DIR, '铝电解槽参数优化_项目分析总览.pptx')
    prs.save(output_path)
    print(f'PPT saved to: {output_path}')
    return output_path


if __name__ == '__main__':
    build_ppt()
